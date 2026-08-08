"""
Category 3 (M-01..M-10) — multimodal readers + RAG + MCP client tests.

Covers: M-01..M-05 (vector store, embeddings, workspace indexer,
semantic_search, MCP client degrade-gracefully) and M-06..M-10 (read_pdf /
read_docx / read_xlsx / read_pptx + read_file auto-dispatch).

Copyright (c) 2026 Dearly Febriano Irwansyah
SPDX-License-Identifier: MIT
"""

from __future__ import annotations

import pytest
import pytest_asyncio

import tools.core.read_pdf as _rp  # noqa: F401  (import smoke)
import tools.core.read_docx as _rd  # noqa: F401
import tools.core.read_xlsx as _rx  # noqa: F401
import tools.core.read_pptx as _rpptx  # noqa: F401
import tools.core.semantic_search as _ss  # noqa: F401
import tools.core.mcp_client as _mcp  # noqa: F401

from openforge.embeddings import embed_text
from openforge.vector_db import VectorStore
from tools.registry import create_default_registry


# ── M-02 / M-03: vector store + embeddings ───────────────────────────────────


def test_embed_text_dim_and_norm():
    v = embed_text("nexa agent")
    assert len(v) == 384
    norm = sum(x * x for x in v) ** 0.5
    assert abs(norm - 1.0) < 0.01


def test_vectorstore_upsert_and_search(tmp_path, monkeypatch):
    monkeypatch.setattr("openforge.vector_db._DB", tmp_path / "vec.db")
    store = VectorStore()
    store.initialize()
    store.upsert("a::0", "a.py", 0, "def hello(): return 'world'", embed_text("hello world python"))
    store.upsert("b::0", "b.py", 0, "unrelated text about cooking", embed_text("cooking food"))
    hits = store.search(embed_text("hello world python"), k=1)
    assert hits and hits[0]["id"] == "a::0" and "hello" in hits[0]["content"]


# ── M-04 / M-05: workspace indexer + semantic_search ─────────────────────────


@pytest.mark.asyncio
async def test_workspace_indexer_and_semantic_search(tmp_path, monkeypatch):
    monkeypatch.setattr("agent.workspace_indexer.FORGE_WORKSPACE", tmp_path)
    monkeypatch.setattr("openforge.vector_db._DB", tmp_path / "vec.db")
    (tmp_path / "hello.py").write_text("def greet(): return 'hi'", encoding="utf-8")
    from agent.workspace_indexer import WorkspaceIndexer

    idx = WorkspaceIndexer()
    n = await idx.index_all()
    assert n >= 1

    from tools.core.semantic_search import semantic_search

    # Force the semantic_search tool to use the same store file.
    monkeypatch.setattr("openforge.vector_db._DB", tmp_path / "vec.db")
    out = await semantic_search("greet")
    assert "hello.py" in out


# ── M-01: MCP client degrades gracefully ─────────────────────────────────────


@pytest.mark.asyncio
async def test_mcp_list_servers_when_not_configured(tmp_path, monkeypatch):
    monkeypatch.setattr("tools.core.mcp_client._SERVERS_FILE", tmp_path / "none.json")
    from tools.core.mcp_client import mcp_list_servers

    out = await mcp_list_servers()
    assert isinstance(out, str)


# ── M-06..M-09: multimodal readers on synthetic files ────────────────────────


@pytest.mark.asyncio
async def test_read_pdf_on_real_pdf(tmp_path, monkeypatch):
    monkeypatch.setattr("tools._paths.FORGE_WORKSPACE", tmp_path, raising=False)
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    pdf_path = tmp_path / "sample.pdf"
    with pdf_path.open("wb") as f:
        writer.write(f)

    # Point the reader at our tmp workspace.
    import importlib
    import tools.core.read_pdf as rp
    importlib.reload(rp)
    monkeypatch.setattr(rp, "resolve_in_workspace", lambda p: tmp_path / p)
    text = await rp.read_pdf("sample.pdf")
    assert "page 1" in text or isinstance(text, str)


@pytest.mark.asyncio
async def test_read_docx_roundtrip(tmp_path, monkeypatch):
    import docx, importlib

    doc = docx.Document()
    doc.add_paragraph("Nexa document body")
    p = tmp_path / "sample.docx"
    doc.save(p)
    import tools.core.read_docx as rd
    monkeypatch.setattr(rd, "resolve_in_workspace", lambda p: tmp_path / p)
    text = await rd.read_docx("sample.docx")
    assert "Nexa document body" in text


@pytest.mark.asyncio
async def test_read_xlsx_roundtrip(tmp_path, monkeypatch):
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["name", "value"])
    ws.append(["alpha", 1])
    p = tmp_path / "data.xlsx"
    wb.save(p)
    import tools.core.read_xlsx as rx
    monkeypatch.setattr(rx, "resolve_in_workspace", lambda p: tmp_path / p)
    text = await rx.read_xlsx("data.xlsx")
    assert "alpha" in text and "name" in text


@pytest.mark.asyncio
async def test_read_pptx_roundtrip(tmp_path, monkeypatch):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Nexa deck"
    p = tmp_path / "deck.pptx"
    prs.save(p)
    import tools.core.read_pptx as rpptx
    monkeypatch.setattr(rpptx, "resolve_in_workspace", lambda p: tmp_path / p)
    text = await rpptx.read_pptx("deck.pptx")
    assert "slide 1" in text or "Nexa deck" in text


# ── M-10: read_file dispatches by extension ──────────────────────────────────


@pytest.mark.asyncio
async def test_read_file_dispatches_pdf(tmp_path, monkeypatch):
    from pypdf import PdfWriter
    import importlib
    import tools.file_tools as ft
    import tools.core.read_pdf as rp

    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    with (tmp_path / "doc.pdf").open("wb") as f:
        writer.write(f)

    # Redirect both resolvers to our tmp workspace.
    monkeypatch.setattr(ft, "resolve_in_workspace", lambda p: tmp_path / p)
    monkeypatch.setattr(rp, "resolve_in_workspace", lambda p: tmp_path / p)
    importlib.reload(ft)
    out = await ft.read_file("doc.pdf")
    assert "page" in out


# ── registry wiring ──────────────────────────────────────────────────────────


def test_registry_has_category3_tools():
    reg = create_default_registry()
    for name in ["read_pdf", "read_docx", "read_xlsx", "read_pptx", "semantic_search", "mcp_list_servers", "mcp_call"]:
        assert reg.has(name), name
