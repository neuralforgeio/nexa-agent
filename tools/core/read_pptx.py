"""Read a .pptx presentation: slide title, body text, speaker notes."""
from __future__ import annotations

from typing import Any

from tools._paths import resolve_in_workspace

try:
    from pptx import Presentation  # python-pptx
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore[assignment]

READ_PPTX_SCHEMA = {
    "type": "object",
    "properties": {
        "path": {"type": "string", "description": "Workspace-relative path to the .pptx file."},
    },
    "required": ["path"],
}


async def read_pptx(path: str, **_: Any) -> str:
    """Extract each slide's text plus speaker notes."""
    if Presentation is None:
        return "read_pptx requires the 'python-pptx' package (pip install python-pptx)."
    full = resolve_in_workspace(path)
    if not full.exists():
        raise ValueError(f"file not found: '{path}'")
    prs = Presentation(str(full))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"=== slide {i} ===")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                out.append(shape.text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            out.append(f"[notes] {slide.notes_slide.notes_text_frame.text}")
    return "\n".join(out) or "(empty presentation)"
